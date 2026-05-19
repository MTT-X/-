
"""Generate presentation.pptx from ShanghaiTech template."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

cwd = os.getcwd()
tmpl_path = os.path.join(cwd, '_tmpl.pptx')
out_path = os.path.join(cwd, 'presentation.pptx')

prs = Presentation(tmpl_path)

# Delete all placeholder slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

LAYOUT_TITLE = prs.slide_layouts[19]    # 1_title slide
LAYOUT_SECTION = prs.slide_layouts[2]    # section header
LAYOUT_CONTENT = prs.slide_layouts[8]    # 3_with content
LAYOUT_TITLE_ONLY = prs.slide_layouts[3] # title only

DARK = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xB9, 0x8D, 0x4B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)


def add_slide(layout):
    return prs.slides.add_slide(layout)


def set_title(slide, text, sz=Pt(30)):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ''
            p = ph.text_frame.paragraphs[0]
            p.text = text
            for run in p.runs:
                run.font.size = sz
                run.font.bold = True
                run.font.color.rgb = DARK
            return ph
    return None


def tb(slide, l, t, w, h, text, sz=Pt(14), bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    """Add textbox"""
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = sz
        run.font.bold = bold
        run.font.color.rgb = color
    return tf


def bullets(slide, l, t, w, h, items, sz=Pt(13)):
    """Add bulleted list"""
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.size = sz
            run.font.color.rgb = BLACK
    return tf


def tbl(slide, l, t, w, h, headers, rows, widths=None):
    """Add table"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    ts = slide.shapes.add_table(n_rows, n_cols, Emu(l), Emu(t), Emu(w), Emu(h))
    table = ts.table
    if widths:
        for i, wd in enumerate(widths):
            table.columns[i].width = Emu(wd)
    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        p.text = hdr
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = BLACK
    return table


# ===== SLIDE 1: COVER =====
s1 = add_slide(LAYOUT_TITLE)
tb(s1, Cm(3), Cm(4), Cm(19), Cm(3),
   'Shipping Insurance Intelligent Classification\nand Risk Assessment',
   Pt(34), True, DARK, PP_ALIGN.CENTER)
tb(s1, Cm(3), Cm(9), Cm(19), Cm(2),
   'FINA2003 Financial Technology Frontiers - Final Project',
   Pt(16), False, GRAY, PP_ALIGN.CENTER)
tb(s1, Cm(3), Cm(12), Cm(19), Cm(2),
   'Group 1: Hua Yu / Mei Yuyang / Yang Yuxuan\n'
   'Partner: Anlan Information Technology  |  May 20, 2025',
   Pt(13), False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 2: OUTLINE =====
s2 = add_slide(LAYOUT_CONTENT)
set_title(s2, 'Outline', Pt(28))
bullets(s2, Cm(2.5), Cm(3.5), Cm(9), Cm(10), [
    'Module 1: Industry & Product Analysis',
    '   - Market landscape and pain points',
    '   - Competitive research',
    '   - Product design',
], Pt(14))
bullets(s2, Cm(12.5), Cm(3.5), Cm(9), Cm(10), [
    'Module 2: Data Processing & Classification',
    '   - Data exploration and preprocessing',
    '   - Classification implementation',
    '   - Error analysis',
    'Module 3: Risk Prediction Exploration',
], Pt(14))

# ===== SLIDE 3: MARKET =====
s3 = add_slide(LAYOUT_CONTENT)
set_title(s3, '1.1 Shipping Insurance: Market & Pain Points', Pt(26))
tb(s3, Cm(2), Cm(3.2), Cm(10), Cm(0.8), 'Market Overview', Pt(16), True, DARK)
bullets(s3, Cm(2), Cm(4), Cm(10), Cm(6), [
    'China shipping insurance: > RMB 50 billion',
    'Annual growth: 8-10%',
    'Key players: Ping An, PICC, CPIC',
    'Tech-driven transformation is industry trend',
], Pt(12))
tb(s3, Cm(12.5), Cm(3.2), Cm(10), Cm(0.8), 'Core Pain Points', Pt(16), True, DARK)
bullets(s3, Cm(12.5), Cm(4), Cm(10), Cm(6), [
    'Free-text descriptions, no standard format',
    'Heavy manual work: 5-15 min per case',
    'Inconsistent classification across staff',
    'Leads to pricing errors & disputes',
], Pt(12))
tb(s3, Cm(2), Cm(11.5), Cm(22), Cm(1.5),
   'Key Insight: 28,304 unstructured descriptions represent real efficiency loss and pricing risk.',
   Pt(11), False, GRAY)

# ===== SLIDE 4: COMPETITIVE =====
s4 = add_slide(LAYOUT_CONTENT)
set_title(s4, '1.2 Competitive Landscape', Pt(26))
tbl(s4, Cm(2), Cm(3.5), Cm(21), Cm(6),
    ['Solution', 'Representative', 'Advantages', 'Disadvantages'],
    [
        ['Rule Engine', 'Insurer in-house', 'Controllable, explainable', 'High maintenance'],
        ['ML Classification', 'Concirrus, Planck', 'Adaptive, high accuracy', 'Needs labeled data'],
        ['LLM-based', 'Zest AI, pilots', 'Semantic understanding', 'Costly, hallucinations'],
        ['Hybrid', 'Select MGAs', 'Rule + ML combo', 'High complexity'],
    ],
    [Cm(4), Cm(5), Cm(6), Cm(6)])
tb(s4, Cm(2), Cm(10.5), Cm(21), Cm(0.8), 'Our Choice: Rule-Based Approach', Pt(15), True, DARK)
bullets(s4, Cm(2), Cm(11.5), Cm(21), Cm(2.5), [
    'Rationale: No labeled data needed, fully interpretable, flexible rules, rapid deployment',
    'Limitations addressable via iterative improvement (TF-IDF fallback, future ML integration)',
], Pt(12))

# ===== SLIDE 5: PRODUCT DESIGN =====
s5 = add_slide(LAYOUT_CONTENT)
set_title(s5, '1.3 Product Design: System Architecture', Pt(26))
tb(s5, Cm(2), Cm(3), Cm(21), Cm(0.8), 'Three-Tier Classification Pipeline:', Pt(15), True, DARK)
bullets(s5, Cm(2), Cm(4.2), Cm(21), Cm(3), [
    'Tier 1: Text Preprocessing - newline removal, fullwidth-to-halfwidth, whitespace normalization',
    'Tier 2: Exclusion Screening - match 10 excluded categories; hit = "High Risk (Decline)"',
    'Tier 3: HS Category Matching - match 22 HS categories via keyword scoring',
], Pt(13))
tb(s5, Cm(2), Cm(8), Cm(21), Cm(0.8), 'Human Review Mechanism:', Pt(15), True, DARK)
bullets(s5, Cm(2), Cm(9.2), Cm(21), Cm(2.5), [
    'Only "High Risk" + "Unclassified" (~30%) trigger manual review',
    '~70% efficiency gain; manual decisions feed back to keyword dictionary',
], Pt(13))
tb(s5, Cm(2), Cm(12.5), Cm(21), Cm(1),
   'Conservative strategy: better to over-reject than under-reject during early deployment.',
   Pt(10), False, GRAY)

# ===== SLIDE 6: SECTION - MODULE 2 =====
s6 = add_slide(LAYOUT_SECTION)
set_title(s6, 'Module 2', Pt(36))
tb(s6, Cm(2.5), Cm(8), Cm(20), Cm(2),
   'Data Processing & Classification Modeling',
   Pt(22), False, WHITE, PP_ALIGN.CENTER)
tb(s6, Cm(2.5), Cm(10.5), Cm(20), Cm(1.5),
   'From 28,304 unstructured text entries to automated classification labels',
   Pt(14), False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 7: DATA EXPLORATION =====
s7 = add_slide(LAYOUT_CONTENT)
set_title(s7, '2.1 Data Exploration', Pt(28))
tbl(s7, Cm(2), Cm(3.5), Cm(21), Cm(3.5),
    ['Dataset', 'Size', 'Content'],
    [
        ['Cargo Description 2025', '28,304 entries', 'Free-text cargo descriptions from policyholders'],
        ['Customs HS Code Table', '8,045 entries', '22 major + detailed HS categories'],
        ['Insurance Exclusion List', '10 categories', 'Goods excluded from coverage'],
    ],
    [Cm(7), Cm(4), Cm(10)])
tb(s7, Cm(2), Cm(7.8), Cm(21), Cm(0.8), 'Data Characteristics:', Pt(14), True, DARK)
bullets(s7, Cm(2), Cm(8.8), Cm(21), Cm(5), [
    '~60% English, ~40% Chinese with mixed usage; avg length ~45 chars',
    '~0.9% (259 entries) invalid: "SEE ATTACHMENT", single special chars',
    'Common noise: container numbers, shipping info, HS codes mixed in text',
    'Short texts dominate (min 1 char), few exceed 500 chars (multi-product)',
], Pt(12))

# ===== SLIDE 8: METHOD =====
s8 = add_slide(LAYOUT_CONTENT)
set_title(s8, '2.2-2.3 Preprocessing & Classification Method', Pt(24))
tb(s8, Cm(1.5), Cm(3), Cm(11), Cm(0.8), 'Preprocessing Pipeline', Pt(14), True, DARK)
bullets(s8, Cm(1.5), Cm(4), Cm(11), Cm(4), [
    '1) Newline normalization: \\n, \\t -> space',
    '2) Fullwidth-to-halfwidth conversion',
    '3) Multi-space collapse',
    'Goal: preserve semantics, avoid aggressive tokenization',
], Pt(11))
tb(s8, Cm(13), Cm(3), Cm(11), Cm(0.8), 'Keyword Dictionary', Pt(14), True, DARK)
bullets(s8, Cm(13), Cm(4), Cm(11), Cm(4), [
    '10 exclusion + 22 HS categories',
    '~500+ Chinese & English keywords',
    'Chinese: substring matching',
    'English short words: regex word-boundary',
], Pt(11))
tb(s8, Cm(1.5), Cm(8.5), Cm(22.5), Cm(0.8),
   'KEY FINDING: Substring False-Match Problem', Pt(14), True, DARK)
bullets(s8, Cm(1.5), Cm(9.3), Cm(22.5), Cm(4), [
    '95.3% of initial exclusion cases hit only 1 keyword - massive false positives',
    '"car" matched "cardboard"; "art" matched "parts"; "ship" matched "shipper"',
    'Fix: word-boundary regex. Exclusion rate: 28.2% -> 17.5% (3,045 false positives removed)',
    'Insight: precision, not recall, is the real challenge in rule-based NLP systems',
], Pt(11))

# ===== SLIDE 9: RESULTS =====
s9 = add_slide(LAYOUT_CONTENT)
set_title(s9, '2.4 Classification Results', Pt(28))
tb(s9, Cm(2), Cm(3.2), Cm(21), Cm(1),
   'Total: 28,304  |  Insurable: 20,309 (82.5%)  |  Excluded: 4,950 (17.5%)  |  Coverage: 83.3%',
   Pt(13), True, DARK, PP_ALIGN.CENTER)
tbl(s9, Cm(1.5), Cm(4.8), Cm(11), Cm(5.5),
    ['Top 5 Exclusion Reasons', 'Count', '%'],
    [
        ['Glass/Ceramic/Stone/Fragile', '1,470', '29.7%'],
        ['Vehicles/Aircraft/Vessels', '1,156', '23.4%'],
        ['Perishable/Fresh Goods', '630', '12.7%'],
        ['Precision Instruments/Chips', '409', '8.3%'],
        ['Flammable/Explosive', '295', '6.0%'],
    ], [Cm(7.5), Cm(2), Cm(1.5)])
tbl(s9, Cm(13), Cm(4.8), Cm(11), Cm(5.5),
    ['Top 5 HS Categories (Insurable)', 'Count', '%'],
    [
        ['Machinery/Electrical Equipment', '3,774', '16.2%'],
        ['Base Metals & Articles', '3,462', '14.8%'],
        ['Plastics & Rubber', '2,389', '10.2%'],
        ['Chemical Products', '1,756', '7.5%'],
        ['Textiles & Articles', '1,552', '6.6%'],
    ], [Cm(7.5), Cm(2), Cm(1.5)])
tb(s9, Cm(2), Cm(11), Cm(21), Cm(1.5),
   'Unclassified: 4,718 (20.2%)  |  Data Quality Issues: 259 (1.1%)',
   Pt(11), False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 10: ERROR ANALYSIS =====
s10 = add_slide(LAYOUT_CONTENT)
set_title(s10, '2.5 Error Analysis', Pt(26))
tb(s10, Cm(2), Cm(3.2), Cm(21), Cm(0.8), 'False Positive Discovery (Fixed)', Pt(15), True, DARK)
bullets(s10, Cm(2), Cm(4.2), Cm(21), Cm(3), [
    'Root cause: short English keyword substring matching (art->parts, car->cardboard, ship->shipper)',
    'Fix: word-boundary regex. Exclusion rate corrected: 28.2% -> 17.5%',
], Pt(12))
tb(s10, Cm(2), Cm(7.5), Cm(21), Cm(0.8), 'Remaining Unclassified (20.2%)', Pt(15), True, DARK)
bullets(s10, Cm(2), Cm(8.5), Cm(21), Cm(4), [
    'Highly specialized terms: "PHAFFIA RHODOZYMA", "PARA ARAMID FIBER"',
    'Vague descriptions: "HOUSEHOLD ITEMS", "SEE ATTACHMENT"',
    'Multi-product combined descriptions in single entry',
    'Future: TF-IDF + LLM zero-shot for specialized terms',
], Pt(12))

# ===== SLIDE 11: SECTION - MODULE 3 =====
s11 = add_slide(LAYOUT_SECTION)
set_title(s11, 'Module 3', Pt(36))
tb(s11, Cm(2.5), Cm(8), Cm(20), Cm(2),
   'Risk Prediction Exploration',
   Pt(22), False, WHITE, PP_ALIGN.CENTER)
tb(s11, Cm(2.5), Cm(10.5), Cm(20), Cm(1.5),
   'Integrating classification labels into risk assessment',
   Pt(14), False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 12: RISK =====
s12 = add_slide(LAYOUT_CONTENT)
set_title(s12, '3.1-3.2 Risk Stratification & Business Value', Pt(24))
tbl(s12, Cm(2), Cm(3.5), Cm(21), Cm(5.5),
    ['Risk Level', 'Score', 'Count', '%', 'Action'],
    [
        ['High (Decline)', '1.0', '4,950', '17.5%', 'Auto-decline + human review'],
        ['Medium-High', '0.5-0.9', '~1,800', '6.4%', 'Manual pricing'],
        ['Low', '0-0.5', '~16,836', '59.5%', 'Auto-underwrite'],
        ['Unclassified', 'N/A', '4,718', '16.7%', 'Human determination'],
    ],
    [Cm(4.5), Cm(2.5), Cm(2.5), Cm(2), Cm(9.5)])
tb(s12, Cm(2), Cm(9.8), Cm(21), Cm(0.8), 'Business Value:', Pt(14), True, DARK)
bullets(s12, Cm(2), Cm(10.8), Cm(21), Cm(2), [
    '~77% of routine cases covered by automatic system, manual review compressed to ~23%',
    'Thousands of hours saved annually per insurer processing 100K+ policies',
], Pt(12))
tb(s12, Cm(2), Cm(13), Cm(21), Cm(0.8),
   'Future: Add claims history + route risk features via LightGBM/XGBoost for continuous risk scoring.',
   Pt(10), False, GRAY)

# ===== SLIDE 13: SUMMARY =====
s13 = add_slide(LAYOUT_CONTENT)
set_title(s13, '4. Summary & Outlook', Pt(28))
tb(s13, Cm(2), Cm(3), Cm(10.5), Cm(0.8), 'Achievements', Pt(15), True, DARK)
bullets(s13, Cm(2), Cm(4), Cm(10.5), Cm(6), [
    '10 exclusion + 22 HS category dictionary',
    'Three-tier rule engine (~400 lines, ~15s)',
    '83.3% classification coverage',
    'Discovered & fixed substring false-match',
], Pt(12))
tb(s13, Cm(13.5), Cm(3), Cm(10.5), Cm(0.8), 'Future Directions', Pt(15), True, DARK)
bullets(s13, Cm(13.5), Cm(4), Cm(10.5), Cm(6), [
    'Short: TF-IDF vector similarity fallback',
    'Mid: XGBoost with accumulated labeled data',
    'Long: BERT semantic matching + claims data',
], Pt(12))
tb(s13, Cm(2), Cm(11.5), Cm(21), Cm(1.5),
   'Core Value: Unstructured text -> Structured risk labels.\n'
   'Machines handle rules. Humans focus on judgment.',
   Pt(14), True, DARK, PP_ALIGN.CENTER)

# ===== SLIDE 14: THANK YOU =====
s14 = add_slide(LAYOUT_SECTION)
set_title(s14, 'Thank You', Pt(40))
tb(s14, Cm(2.5), Cm(8), Cm(20), Cm(2),
   'Questions & Discussion',
   Pt(22), False, WHITE, PP_ALIGN.CENTER)
tb(s14, Cm(2.5), Cm(10.5), Cm(20), Cm(1.5),
   'Group 1: Hua Yu / Mei Yuyang / Yang Yuxuan\n'
   'FINA2003 Financial Technology Frontiers',
   Pt(14), False, GRAY, PP_ALIGN.CENTER)

# ===== SAVE =====
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Total slides: {len(prs.slides)}')
os.remove(tmpl_path)
print('Done.')
